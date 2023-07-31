from flask import Flask, render_template, request, redirect, url_for, session, send_file
from flask_mysqldb import MySQL
import MySQLdb.cursors
import re
from PyPDF2 import PdfReader, PdfWriter
import os
import secrets
from flask import session
from PIL import Image, ImageDraw, ImageFont
import io
import random
import string
import zipfile
from io import BytesIO
import pandas as pd
from pptx import Presentation
import tempfile
import shutil
import datetime
from datetime import datetime as dt
from collections import defaultdict
import json
from werkzeug.utils import secure_filename
import xlrd
import xlwt
import openpyxl
import subprocess
import pyclamd
import time
from flask import jsonify


app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'


app.secret_key = 'xyzsdfg'

app.config['MYSQL_HOST'] = 'localhost'
app.config['MYSQL_USER'] = 'root'
app.config['MYSQL_PASSWORD'] = ''
app.config['MYSQL_DB'] = 'cdacmeity'

mysql = MySQL(app)


@app.route('/')
@app.route('/homepage')
def homepage():
    return render_template('homepage.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    message = ''
    if request.method == 'POST' and 'email' in request.form and 'password' in request.form and 'captcha' in request.form:
        email = request.form['email']
        password = request.form['password']
        captcha = request.form['captcha']

        # Verify CAPTCHA
        if captcha != session.get('captcha_solution'):
            message = 'CAPTCHA verification failed!'
            return render_template('login.html', message=message)

        cursor = mysql.connection.cursor(MySQLdb.cursors.DictCursor)
        cursor.execute('SELECT * FROM users WHERE email = %s AND password = %s', (email, password,))
        user = cursor.fetchone()
        if user:
            session['loggedin'] = True
            session['uid'] = user['uid']
            session['first_name'] = user['first_name']
            session['last_name'] = user['last_name']
            session['email'] = user['email']
            message = 'Logged in successfully!'
            #return render_template('user_dashboard.html', message=message)
            uploaded_files()
            return redirect(url_for('user_dashboard'))
        else:
            message = 'Please enter correct email/password!'

    return render_template('login.html', message=message)

@app.route('/adminlogin', methods=['GET', 'POST'])
def adminlogin():
    return render_template('admin_login.html')

@app.route('/loginadmin', methods=['GET', 'POST'])
def loginadmin():
    message = ''
    if request.method == 'POST' and 'email' in request.form and 'password' in request.form and 'captcha' in request.form:
        email = request.form['email']
        password = request.form['password']
        captcha = request.form['captcha']

        # Verify CAPTCHA
        if captcha != session.get('captcha_solution'):
            message = 'CAPTCHA verification failed!'
            return render_template('admin_login.html', message=message)
        
        if email == "admin@example.com" and password == "admin123":
            # Admin login success, redirect to admin page
            session['loggedin'] = True
            return redirect(url_for('admin_dashboard'))
        else:
            # Regular user login, show the login page with an error message
            message = "Invalid credentials. Please try again."
            return render_template('admin_login.html', message=message)        

    return render_template('admin_login.html', message=message)

@app.route('/captcha')
def captcha():
    # Generate CAPTCHA solution
    solution = generate_captcha_solution()
    # Store the solution in the session
    session['captcha_solution'] = solution

    # Generate and serve the CAPTCHA image
    image = generate_captcha_image(solution)
    image_io = io.BytesIO()
    image.save(image_io, 'PNG')
    image_io.seek(0)
    return send_file(image_io, mimetype='image/png')


def generate_captcha_solution():
    # Generate a random solution with a combination of uppercase letters and digits
    solution = ''.join(random.choices(string.ascii_uppercase + string.digits, k=4))
    return solution


def generate_captcha_image(solution):
    # Create a new image using PIL
    image = Image.new('RGB', (200, 100), (255, 255, 255))
    draw = ImageDraw.Draw(image)

    # Load a font
    font = ImageFont.truetype('arial.ttf', 40)

    # Add noise to the background
    for _ in range(1000):
        x = random.randint(0, 199)
        y = random.randint(0, 99)
        draw.point((x, y), fill=random_color())

    # Draw the solution on the image with random colors and positions
    for i, char in enumerate(solution):
        x = random.randint(10 + i * 40, 30 + i * 40)
        y = random.randint(10, 40)
        draw.text((x, y), char, font=font, fill=random_color())

    return image


def random_color():
    # Generate a random RGB color tuple
    return tuple(random.choices(range(256), k=3))

@app.route('/forgotpassword', methods=['GET', 'POST'])
def forgotpassword():
    message = ''
    if request.method == 'POST' and 'email' in request.form:
        email = request.form['email']

        cursor = mysql.connection.cursor(MySQLdb.cursors.DictCursor)
        cursor.execute('SELECT * FROM users WHERE email = %s', (email,))
        user = cursor.fetchone()

        if user:
            return redirect(url_for('/resetpassword', email=email))
        else:
            message = 'Please enter correct email.'

    return render_template('forgot_password.html', message=message)

@app.route('/resetpassword', methods=['GET', 'POST'])
def resetpassword():
    if 'loggedin' in session:
        user_id = session['uid']
        email = session['email']
        message = ''

        if not email:
            return render_template('reset_password.html')

        if request.method == 'POST' and 'password' in request.form and 'confirm_password' in request.form:
            password = request.form['password']
            confirm_password = request.form['confirm_password']

            if not password:
                message = 'Please fill out the password!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            if not confirm_password:
                message = 'Please fill out the confirm password!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)
            
            if not re.search(r'^(?=.*[!@#$%^&*])[\w!@#$%^&*]{8,15}$', password):
                message = 'Password must be between 8 to 15 characters long and contain at least one special character (!@#$%^&*)'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            if password != confirm_password:
                message = 'Password and confirm password do not match!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            # Update the user's password in the database
            cursor = mysql.connection.cursor()
            cursor.execute('UPDATE users SET password = %s WHERE email = %s' , (password, email))
            cursor.execute('UPDATE users SET confirm_password = %s WHERE email = %s', (confirm_password, email))
            mysql.connection.commit()

            message = 'Your password has been reset successfully.'

        return render_template('reset_password.html', email=email, message=message)
    return render_template('forgot_password.html', email=email, message=message)

@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    user_id = session['uid']
    message = ''
    if request.method == 'POST' and 'email' in request.form:
        email = request.form['email']

        cursor = mysql.connection.cursor(MySQLdb.cursors.DictCursor)
        cursor.execute('SELECT * FROM users WHERE email = %s AND uid = %s', (email, user_id))
        user = cursor.fetchone()

        if user:
            return redirect(url_for('reset_password', email=email))
        else:
            message = 'Please enter correct email.'

    return render_template('forgot_password.html', message=message)


@app.route('/reset-password', methods=['GET', 'POST'])
def reset_password():
    if 'loggedin' in session:
        user_id = session['uid']
        email = session['email']
        message = ''

        if not email:
            return render_template('reset_password.html')

        if request.method == 'POST' and 'password' in request.form and 'confirm_password' in request.form:
            password = request.form['password']
            confirm_password = request.form['confirm_password']

            if not password:
                message = 'Please fill out the password!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            if not confirm_password:
                message = 'Please fill out the confirm password!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)
            
            if not re.search(r'^(?=.*[!@#$%^&*])[\w!@#$%^&*]{8,15}$', password):
                message = 'Password must be between 8 to 15 characters long and contain at least one special character (!@#$%^&*)'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            if password != confirm_password:
                message = 'Password and confirm password do not match!'
                form_data = request.form
                return render_template('reset_password.html', message=message, form_data=form_data)

            # Update the user's password in the database
            cursor = mysql.connection.cursor()
            cursor.execute('UPDATE users SET password = %s WHERE email = %s' , (password, email))
            cursor.execute('UPDATE users SET confirm_password = %s WHERE email = %s', (confirm_password, email))
            mysql.connection.commit()

            message = 'Your password has been reset successfully.'

        return render_template('reset_password.html', email=email, message=message)
    return render_template('forgot_password.html', email=email, message=message)

@app.route('/logout')
def logout():
    session.pop('loggedin', None)
    session.pop('uid', None)
    session.pop('email', None)
    return redirect(url_for('login'))

@app.route('/adminlogout')
def adminlogout():
    session.pop('loggedin', None)
    session.pop('email', None)
    return redirect(url_for('adminlogin'))


@app.route('/register', methods=['GET', 'POST'])
def register():
    message = ''
    form_data = {}
    if request.method == 'POST':
        salutation = request.form['salutation']
        first_name = request.form['first_name']
        middle_name = request.form['middle_name']
        last_name = request.form['last_name']
        email = request.form['email']
        password = request.form['password']
        confirm_password = request.form['confirm_password']
        employee_id=request.form['employee_id']
        group_name=request.form['group_name']

        cursor = mysql.connection.cursor(MySQLdb.cursors.DictCursor)
        cursor.execute('SELECT * FROM users WHERE email = %s', (email,))
        account = cursor.fetchone()
        if account:
            message = 'Account already exists!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not salutation:
            message = 'Please fill out the salutation!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)
        
        if not employee_id:
            message = 'Please fill out the Employee ID field'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)
        
        if not group_name:
            message = 'Please fill out the Group Name!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not first_name:
            message = 'Please fill out the first name!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not last_name:
            message = 'Please fill out the last name!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not email:
            message = 'Please fill out the email!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not password:
            message = 'Please fill out the password!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if not confirm_password:
            message = 'Please fill out the confirm password!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)
        
        if len(employee_id)!=6:
            message = 'Employee ID must be 6 digits long!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)
        
        if not re.search(r'^(?=.*[!@#$%^&*])[\w!@#$%^&*]{8,15}$', password):
            message = 'Password must be between 8 to 15 characters long and contain at least one special character (!@#$%^&*)'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        if password != confirm_password:
            message = 'Password and confirm password do not match!'
            form_data = request.form
            return render_template('register.html', message=message, form_data=form_data)

        cursor.execute('INSERT INTO users VALUES (NULL, %s, %s, %s, %s, %s, %s, %s, %s, %s, 0, 0, 0, 0, 0, 0)', (salutation, first_name, middle_name, last_name, email, employee_id, group_name, password, confirm_password))
        mysql.connection.commit()
        message = 'You have successfully registered, please login :)'
        return render_template('login.html', message=message, form_data=form_data)
    return render_template('register.html', message=message, form_data=form_data)


@app.route('/admin_dashboard')
def admin_dashboard():
    if 'loggedin' in session:

        #Query the database to retrive users
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT * FROM users")
        users = cursor.fetchall()

        # Count the total number of users
        total_users = len(users)

        # Query the database to retrieve the user's uploaded files
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT * FROM files")
        uploaded_files = cursor.fetchall()
        cursor.close()

        # Print the uploaded_files list for debugging
        print(uploaded_files)

        # Count the total number of files
        total_files = len(uploaded_files)
        print(total_files)

        # Count the number of files with PDF extension
        pdf_files = sum(file[3].lower().endswith('.pdf') for file in uploaded_files)
        print(pdf_files)

        # Count the number of files with PPT extension
        ppt_files = sum(file[3].lower().endswith('.ppt') or file[3].lower().endswith('.pptx') for file in uploaded_files)

        # Count the number of images (extensions: .jpg, .jpeg, .png, .gif)
        image_files = sum(file[3].lower().endswith(('.jpg', '.jpeg', '.png', '.gif')) for file in uploaded_files)

        # Count the number of documents (extensions: .doc, .docx, .txt)
        doc_files = sum(file[3].lower().endswith(('.doc', '.docx', '.txt')) for file in uploaded_files)

        # Count the number of Excel sheets (extensions: .xls, .xlsx)
        xls_files = sum(file[3].lower().endswith(('.xls', '.xlsx')) for file in uploaded_files)

        return render_template('admin_dashboard.html', total_users=total_users, users= users, files=uploaded_files, total_files=total_files, pdf_files = pdf_files, ppt_files = ppt_files, image_files = image_files, 
                               doc_files = doc_files, xls_files = xls_files,)
    return redirect(url_for('adminlogin'))

@app.route('/view', methods = ['GET', 'POST'])
def view():
    if 'loggedin' in session:
        user_id = request.args.get('userid')

        # Query the database to retrieve the user's uploaded files
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT * FROM files WHERE user_id = %s", (user_id,))
        uploaded_files = cursor.fetchall()
        cursor.close()

        # Query the database to retrieve the user's first name and last name
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT first_name, last_name FROM users WHERE uid = %s", (user_id,))
        user_data = cursor.fetchone()
        cursor.close()

        first_name = user_data[0]
        last_name = user_data[1]

        # Print the uploaded_files list for debugging
        print(uploaded_files)

        # Count the total number of files
        total_files = len(uploaded_files)
        print(total_files)

        # Count the number of files with PDF extension
        pdf_files = sum(file[3].lower().endswith('.pdf') for file in uploaded_files)
        print(pdf_files)

        # Count the number of files with PPT extension
        ppt_files = sum(file[3].lower().endswith('.ppt') or file[3].lower().endswith('.pptx') for file in uploaded_files)

        # Count the number of images (extensions: .jpg, .jpeg, .png, .gif)
        image_files = sum(file[3].lower().endswith(('.jpg', '.jpeg', '.png', '.gif')) for file in uploaded_files)

        # Count the number of documents (extensions: .doc, .docx, .txt)
        doc_files = sum(file[3].lower().endswith(('.doc', '.docx', '.txt')) for file in uploaded_files)

        # Count the number of Excel sheets (extensions: .xls, .xlsx)
        xls_files = sum(file[3].lower().endswith(('.xls', '.xlsx')) for file in uploaded_files)

        # Count the number of files uploaded per month
        current_year = dt.now().year
        files_per_month = [0] * 12  # Initialize a list with 12 elements, representing the months

        for file in uploaded_files:
            file_month = file[4].month
            files_per_month[file_month - 1] += 1  # Increment the count for the corresponding month

        print("*****************")
        print(files_per_month)  # Display the files per month on the console
        
        return render_template('user_view.html', files=uploaded_files, total_files=total_files, pdf_files = pdf_files, ppt_files = ppt_files, image_files = image_files, 
                               doc_files = doc_files, xls_files = xls_files, files_per_month=files_per_month, current_year=current_year,
                               first_name=first_name, last_name=last_name)
    return redirect(url_for('adminlogin'))

@app.route('/deleteuser')
def deleteuser():
    if 'loggedin' in session:
        user_id = request.args.get('userid')
        # Query the database to delete the user's data
        cursor = mysql.connection.cursor()
        cursor.execute("DELETE FROM users WHERE uid = %s", (user_id,))
        mysql.connection.commit()
        cursor.close()

        # Query the database to delete the user's uploaded files
        cursor = mysql.connection.cursor()
        cursor.execute("DELETE FROM files WHERE user_id = %s", (user_id,))
        mysql.connection.commit()
        cursor.close()

        return redirect(url_for('admin_dashboard'))
    return redirect(url_for('adminlogin'))

@app.route('/user_dashboard')
def user_dashboard():
    if 'loggedin' in session:
        user_id = session['uid']
        # Query the database to retrieve the user's uploaded files
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT * FROM files WHERE user_id = %s", (user_id,))
        uploaded_files = cursor.fetchall()
        cursor.close()

        # Print the uploaded_files list for debugging
        print(uploaded_files)

        # Count the total number of files
        total_files = len(uploaded_files)
        print(total_files)

        # Count the number of files with PDF extension
        pdf_files = sum(file[3].lower().endswith('.pdf') for file in uploaded_files)
        print(pdf_files)

        # Count the number of files with PPT extension
        ppt_files = sum(file[3].lower().endswith('.ppt') or file[3].lower().endswith('.pptx') for file in uploaded_files)

        # Count the number of images (extensions: .jpg, .jpeg, .png, .gif)
        image_files = sum(file[3].lower().endswith(('.jpg', '.jpeg', '.png', '.gif')) for file in uploaded_files)

        # Count the number of documents (extensions: .doc, .docx, .txt)
        doc_files = sum(file[3].lower().endswith(('.doc', '.docx', '.txt')) for file in uploaded_files)

        # Count the number of Excel sheets (extensions: .xls, .xlsx)
        xls_files = sum(file[3].lower().endswith(('.xls', '.xlsx')) for file in uploaded_files)

        # Update the counts in the users table
        cursor = mysql.connection.cursor()
        cursor.execute("UPDATE users SET total_files = %s, pdf_files = %s, ppt_files = %s, image_files = %s, doc_files = %s, xls_files = %s WHERE uid = %s",
                       (total_files, pdf_files, ppt_files, image_files, doc_files, xls_files, user_id))
        mysql.connection.commit()
        cursor.close()

        # Count the number of files uploaded per month
        current_year = dt.now().year
        files_per_month = [0] * 12  # Initialize a list with 12 elements, representing the months

        for file in uploaded_files:
            file_month = file[4].month
            files_per_month[file_month - 1] += 1  # Increment the count for the corresponding month

        print("*****************")
        print(files_per_month)  # Display the files per month on the console

        return render_template('user_dashboard.html', first_name=session['first_name'], last_name=session['last_name'],
                               email=session['email'], files=uploaded_files, total_files=total_files, pdf_files = pdf_files, ppt_files = ppt_files, image_files = image_files, 
                               doc_files = doc_files, xls_files = xls_files, files_per_month=files_per_month, current_year=current_year)
    return redirect(url_for('login'))

@app.route('/pdf_user')
def pdf_user():
    if 'loggedin' in session:
        return render_template('user_pdf.html', first_name=session['first_name'], last_name=session['last_name'], email=session['email'])
    return redirect(url_for('login'))

@app.route('/image_user')
def image_user():
    if 'loggedin' in session:
        return render_template('user_image.html', first_name=session['first_name'], last_name=session['last_name'], email=session['email'])
    return redirect(url_for('login'))

@app.route('/docs_user')
def docs_user():
    if 'loggedin' in session:
        return render_template('user_doc.html', first_name=session['first_name'], last_name=session['last_name'], email=session['email'])
    return redirect(url_for('login'))

@app.route('/ppt_user')
def ppt_user():
    if 'loggedin' in session:
        return render_template('user_ppt.html', first_name=session['first_name'], last_name=session['last_name'], email=session['email'])
    return redirect(url_for('login'))

@app.route('/xlsx_user')
def xlsx_user():
    if 'loggedin' in session:
        return render_template('user_xls.html', first_name=session['first_name'], last_name=session['last_name'], email=session['email'])
    return redirect(url_for('login'))

@app.route('/profile')
def profile():
    return render_template('profile.html')

@app.route('/compress', methods=['POST'])
def compress():
    if 'loggedin' in session:
        if 'file' not in request.files:
            return redirect(url_for('pdf_user'))

        file = request.files['file']
        if file.filename == '':
            return redirect(url_for('pdf_user'))

        if file and allowed_file(file.filename):
            # Save the uploaded file to a temporary location
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            # Get the absolute file path
            abs_file_path = os.path.abspath(file_path)

            # Display loading message before starting the scanning process
            message = "Scanning the file, please wait..."
            time.sleep(3)  # Simulating the scanning process (remove this line in your actual code)

            # Process the file_path, perform operations, and return response
            # Example: Perform virus scan on the file
            virus_scan_result = scan_file(abs_file_path)

            # Return the virus scan result or any other response
            if virus_scan_result == 'Virus found':
                message = "Virus Detected in file uploaded."
                # Delete the temporary file after processing, if needed
                os.remove(abs_file_path)
                return render_template('user_pdf.html', message=message)
            elif virus_scan_result == 'File is safe':
                # Compress the PDF file
                compressed_file_path = compress_pdf(file_path)

                # Delete the original uploaded file
                os.remove(file_path)

                # Store file details in the database (you need to implement this function)
                store_file_details(compressed_file_path)

                # Create a response with the compressed file
                response = send_file(compressed_file_path, download_name='compressed.pdf', as_attachment=True)

                # Set the message for file compressed
                message = "File Compressed"

                # Add the message to the response as a cookie
                response.set_cookie('message', message)

                return response
            else:
                message = "Unknown Result"
                return render_template('user_pdf.html', message=message)

    return redirect(url_for('pdf_user'))


#pyclamav use
def scan_file(abs_file_path):
    try:
        cd = pyclamd.ClamdAgnostic()
        cd.ping()  # Check if ClamAV daemon is running
        result = cd.scan_file(abs_file_path)
        print("File Path *********")
        print(abs_file_path)
        print("Result *********")
        print(result)
        if result is None:
            return "File is safe"
        elif result[abs_file_path] == 'OK':
            return "File is safe"
        else:
            return "Virus found"
    except pyclamd.ConnectionError:
        return "Error: Could not connect to ClamAV daemon"
    except Exception as e:
        return "Error: {}".format(str(e))
    
def store_file_details(file_path):
    # Extract file information
    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    # Get the current user's ID
    user_id = session['uid']

    # Get the current date and time
    current_date = datetime.datetime.now()

    # Update the file details in the database
    cursor = mysql.connection.cursor()
    cursor.execute('INSERT INTO files VALUES (NULL, %s, %s, %s, %s)', (user_id, file_name, file_extension, current_date))
    mysql.connection.commit()


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() == 'pdf'


def compress_pdf(file_path):
    output_file_path = f"compressed_{file_path}"
    output_dir = os.path.dirname(output_file_path)
    os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist
    with open(file_path, "rb") as file:
        pdf = PdfReader(file)
        writer = PdfWriter()
        writer.compress = True  # Enable compression
        for page in pdf.pages:
            writer.add_page(page)
        with open(output_file_path, "wb") as output_file:
            writer.write(output_file)
    return output_file_path
    
@app.route('/compress-image', methods=['POST'])
def compress_image():
    if 'loggedin' in session:
        if 'file' not in request.files:
            return redirect(url_for('image_user'))

        file = request.files['file']
        compress_level = request.form.get('compress_level')
        print("******** Compress Level ************ ")
        print(compress_level)
        if file.filename == '':
            return redirect(url_for('image_user'))

        if file and allowed_image_file(file.filename):
            # Save the uploaded file to a temporary location
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            # Get the absolute file path
            abs_file_path = os.path.abspath(file_path)
            print("File Path MAIN ONE *********")
            print(abs_file_path)

            # Process the file_path, perform operations, and return response
            # Example: Perform virus scan on the file
            virus_scan_result = scan_file(abs_file_path)
            
            # Return the virus scan result or any other response
            print(virus_scan_result)
            if virus_scan_result == 'Virus found':
                message = "Virus Detected in file uploaded."
                # Delete the temporary file after processing, if needed
                os.remove(abs_file_path)
                return render_template('user_image.html', message=message)
            elif virus_scan_result == 'File is safe':
               # Compress the image file
                compressed_file_path = compress_image_file(file_path, compress_level)
                response = send_file(compressed_file_path, download_name='compressed.jpg',as_attachment=True)

                # Delete the original uploaded file
                print("FILLLLLEEE PATTTTHHH")
                print(file_path)
                os.remove(file_path)

                # Store file details in the database
                store_file_details(compressed_file_path)
                return response
            else:
                message = "Unknown Result"
                return render_template('user_image.html', message=message)

    return redirect(url_for('image_user'))

#pyclamav use
def scan_file(abs_file_path):
    try:
        cd = pyclamd.ClamdAgnostic()
        cd.ping()  # Check if ClamAV daemon is running
        result = cd.scan_file(abs_file_path)
        print("File Path *********")
        print(abs_file_path)
        print("Result *********")
        print(result)
        if result is None:
            return "File is safe"
        elif result[abs_file_path] == 'OK':
            return "File is safe"
        else:
            return "Virus found"
    except pyclamd.ConnectionError:
        return "Error: Could not connect to ClamAV daemon"
    except Exception as e:
        return "Error: {}".format(str(e))

def store_file_details(file_path):
    # Extract file information
    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    # Get the current user's ID
    user_id = session['uid']

    # Get the current date and time
    current_date = datetime.datetime.now()

    # Update the file details in the database
    cursor = mysql.connection.cursor()
    cursor.execute('INSERT INTO files VALUES (NULL, %s, %s, %s, %s)', (user_id, file_name, file_extension, current_date))
    mysql.connection.commit()

def allowed_image_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ['jpg', 'jpeg', 'png', 'gif']

def compress_image_file(file_path, compress_level):
    output_file_path = f"compressed_{file_path}"
    output_dir = os.path.dirname(output_file_path)
    os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist
    image = Image.open(file_path)
    if compress_level == 'Low':
        image.save(output_file_path, optimize=True, quality=80)
    elif compress_level == 'High':
        image.save(output_file_path, optimize=True, quality=20)
    else:
        image.save(output_file_path, optimize=True, quality=50)
    return output_file_path

@app.route('/compress_document', methods=['POST'])
def compress_document():
    # Get the uploaded file from the request
    file = request.files['file']

    # Save the uploaded file to a temporary location
    temp_file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(temp_file_path)
    # Get the absolute file path
    abs_file_path = os.path.abspath(temp_file_path)
    print("File Path MAIN ONE *********")
    print(abs_file_path)

    # Process the file_path, perform operations, and return response
    # Example: Perform virus scan on the file
    virus_scan_result = scan_file(abs_file_path)
            
    # Return the virus scan result or any other response
    print(virus_scan_result)
    if virus_scan_result == 'Virus found':
        message = "Virus Detected in file uploaded."
        # Delete the temporary file after processing, if needed
        os.remove(abs_file_path)
        return render_template('user_doc.html', message=message)
    elif virus_scan_result == 'File is safe':
        #  Compress the Word document
        compressed_file_path = compress_word_document(temp_file_path)

        # Remove the temporary file
        os.remove(temp_file_path)

        # Store file details in the database
        store_file_details(compressed_file_path)

        # Return the compressed file to the user for download
        return send_file(compressed_file_path, as_attachment=True)
    else:
        message = "Unknown Result"
        return render_template('user_doc.html', message=message)

#pyclamav use
def scan_file(abs_file_path):
    try:
        cd = pyclamd.ClamdAgnostic()
        cd.ping()  # Check if ClamAV daemon is running
        result = cd.scan_file(abs_file_path)
        print("File Path *********")
        print(abs_file_path)
        print("Result *********")
        print(result)
        if result is None:
            return "File is safe"
        elif result[abs_file_path] == 'OK':
            return "File is safe"
        else:
            return "Virus found"
    except pyclamd.ConnectionError:
        return "Error: Could not connect to ClamAV daemon"
    except Exception as e:
        return "Error: {}".format(str(e))

def store_file_details(file_path):
    # Extract file information
    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    # Get the current user's ID
    user_id = session['uid']

    # Get the current date and time
    current_date = datetime.datetime.now()

    # Update the file details in the database
    cursor = mysql.connection.cursor()
    cursor.execute('INSERT INTO files VALUES (NULL, %s, %s, %s, %s)', (user_id, file_name, file_extension, current_date))
    mysql.connection.commit()

def compress_word_document(file_path):
    output_file_path = 'temp/compressed_' + file_path
    output_dir = os.path.dirname(output_file_path)
    os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist

    # Open the Word document as a zip file
    with zipfile.ZipFile(file_path, "r") as zip_file:
        # Create a new in-memory zip file
        compressed_zip = BytesIO()

        with zipfile.ZipFile(compressed_zip, "w", compression=zipfile.ZIP_DEFLATED) as new_zip:
            # Iterate over each file in the original zip
            for item in zip_file.infolist():
                # Read the file data
                data = zip_file.read(item.filename)
                # Compress the file and add it to the new zip
                new_zip.writestr(item, data, compress_type=zipfile.ZIP_DEFLATED)

        # Save the compressed zip to the output file path
        with open(output_file_path, "wb") as output_file:
            output_file.write(compressed_zip.getvalue())

    return output_file_path

@app.route('/compress_excel', methods=['POST'])
def compress_excel():
    # Check if the 'file' field exists in the request
    if 'file' not in request.files:
        return redirect(url_for('xlsx_user'))

    # Get the uploaded file
    file = request.files['file']

    # Check if a file was selected
    if file.filename == '':
        return redirect(url_for('xlsx_user'))

    # Check if the file has an allowed extension (Excel format)
    if file and allowed_excel_file(file.filename):
        # Save the uploaded file to a temporary location
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)

        # Get the absolute file path
        abs_file_path = os.path.abspath(file_path)
        print("File Path MAIN ONE *********")
        print(abs_file_path)

        # Process the file_path, perform operations, and return response
        # Example: Perform virus scan on the file
        virus_scan_result = scan_file(abs_file_path)
            
        # Return the virus scan result or any other response
        print(virus_scan_result)
        if virus_scan_result == 'Virus found':
            message = "Virus Detected in file uploaded."
            # Delete the temporary file after processing, if needed
            os.remove(abs_file_path)
            return render_template('user_xls.html', message=message)
        elif virus_scan_result == 'File is safe':
            # Compress the Excel sheet
            compressed_file_path = compress_excel_sheet(file_path)

            # Store file details in the database
            store_file_details(compressed_file_path)

            # Delete the original uploaded file
            os.remove(file_path)

            # Provide the compressed file as a download
            return send_file(compressed_file_path, download_name='compressed.xlsx')
        else:
            message = "Unknown Result"
            return render_template('user_xls.html', message=message)

    return redirect(url_for('xlsx_user'))

#pyclamav use
def scan_file(abs_file_path):
    try:
        cd = pyclamd.ClamdAgnostic()
        cd.ping()  # Check if ClamAV daemon is running
        result = cd.scan_file(abs_file_path)
        print("File Path *********")
        print(abs_file_path)
        print("Result *********")
        print(result)
        if result is None:
            return "File is safe"
        elif result[abs_file_path] == 'OK':
            return "File is safe"
        else:
            return "Virus found"
    except pyclamd.ConnectionError:
        return "Error: Could not connect to ClamAV daemon"
    except Exception as e:
        return "Error: {}".format(str(e))

def store_file_details(file_path):
    # Extract file information
    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    # Get the current user's ID
    user_id = session['uid']

    # Get the current date and time
    current_date = datetime.datetime.now()

    # Update the file details in the database
    cursor = mysql.connection.cursor()
    cursor.execute('INSERT INTO files VALUES (NULL, %s, %s, %s, %s)', (user_id, file_name, file_extension, current_date))
    mysql.connection.commit()

def allowed_excel_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ['xlsx', 'xls']


def compress_excel_sheet(file_path):
    # Load the Excel file using xlrd
    workbook = xlrd.open_workbook(file_path, formatting_info=True)
    sheets = workbook.sheet_names()

    # Create a new output workbook using xlwt
    compressed_workbook = xlwt.Workbook()

    # Apply optimization techniques based on the compress_level
    for sheet_name in sheets:
        worksheet = workbook.sheet_by_name(sheet_name)
        compressed_worksheet = compressed_workbook.add_sheet(sheet_name)

        apply_high_level_optimization(worksheet, compressed_worksheet)

    # Create a new output file path
    compressed_file_path = f"compressed_{file_path}"

    # Save the modified Excel file using xlwt
    compressed_workbook.save(compressed_file_path)

    return compressed_file_path

def apply_high_level_optimization(worksheet, compressed_worksheet):
    # Apply high-level optimization techniques to the worksheet
    # Example: Simplify formulas, remove redundant data, etc.
    for row in range(worksheet.nrows):
        for col in range(worksheet.ncols):
            cell = worksheet.cell(row, col)
            compressed_worksheet.write(row, col, cell.value)

@app.route('/compress-powerpoint', methods=['POST'])

def compress_powerpoint_route():
    if 'loggedin' in session:
        if 'file' not in request.files:
            return redirect(url_for('ppt_user'))

        file = request.files['file']
        if file.filename == '':
            return redirect(url_for('ppt_user'))

        if file and allowed_powerpoint_file(file.filename):
            # Save the uploaded file to a temporary location
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)
            # Get the absolute file path
            abs_file_path = os.path.abspath(file_path)
            print("File Path MAIN ONE *********")
            print(abs_file_path)

            # Process the file_path, perform operations, and return response
            # Example: Perform virus scan on the file
            virus_scan_result = scan_file(abs_file_path)
                
            # Return the virus scan result or any other response
            print(virus_scan_result)
            if virus_scan_result == 'Virus found':
                message = "Virus Detected in file uploaded."
                # Delete the temporary file after processing, if needed
                os.remove(abs_file_path)
                return render_template('user_ppt.html', message=message)
            elif virus_scan_result == 'File is safe':
                # Compress the PowerPoint presentation
                compressed_file_path = compress_powerpoint_file(file_path)

                # Delete the original uploaded file
                os.remove(file_path)

                # Store file details in the database
                store_file_details(file_path)

                # Provide the compressed file as a download
                return send_file(compressed_file_path, download_name='compressed.pptx')
            else:
                message = "Unknown Result"
                return render_template('user_ppt.html', message=message)

    return redirect(url_for('ppt_user'))

def store_file_details(file_path):
    # Extract file information
    file_name = os.path.basename(file_path)
    file_extension = os.path.splitext(file_name)[1].lower()

    # Get the current user's ID
    user_id = session['uid']

    # Get the current date and time
    current_date = datetime.datetime.now()

    # Update the file details in the database
    cursor = mysql.connection.cursor()
    cursor.execute('INSERT INTO files VALUES (NULL, %s, %s, %s, %s)', (user_id, file_name, file_extension, current_date))
    mysql.connection.commit()

def allowed_powerpoint_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() == 'pptx'

def compress_powerpoint_file(file_path):
    output_file_path = f"compressed_{file_path}"
    output_dir = os.path.dirname(output_file_path)
    os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist

    # Load the PowerPoint presentation
    presentation = Presentation(file_path)

    # Compress the presentation
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # Compress images in the slide
                shape.image.compression = 70  # Set the compression level as desired (0-100)

    # Save the compressed presentation to the output file path
    compressed_file_path = tempfile.mktemp(suffix='.pptx')
    presentation.save(compressed_file_path)

    return compressed_file_path

@app.route('/uploaded_files')
def uploaded_files():
    if 'loggedin' in session:
        user_id = session['uid']
        # Query the database to retrieve the user's uploaded files
        cursor = mysql.connection.cursor()
        cursor.execute("SELECT * FROM files WHERE user_id = %s", (user_id,))
        uploaded_files = cursor.fetchall()
        cursor.close()
        print(uploaded_files)
        return render_template('uploaded_files.html', files=uploaded_files)

if __name__ == '__main__':
    app.run(debug=True)
